from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_NAME = "BOA_Jarvis_to_LeanFT_Migration_Plan.pptx"

BRAND = RGBColor(242, 101, 34)
BRAND_DARK = RGBColor(89, 37, 12)
INK = RGBColor(31, 41, 55)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(30, 64, 175)
GREEN = RGBColor(22, 101, 52)
RED = RGBColor(185, 28, 28)
AMBER = RGBColor(180, 83, 9)
PALE_ORANGE = RGBColor(255, 237, 213)
PALE_BLUE = RGBColor(219, 234, 254)
PALE_GREEN = RGBColor(220, 252, 231)
PALE_RED = RGBColor(254, 226, 226)
PALE_GOLD = RGBColor(254, 249, 195)
PALE_GRAY = RGBColor(241, 245, 249)


def build_deck():
	presentation = Presentation()
	presentation.slide_width = Inches(13.333)
	presentation.slide_height = Inches(7.5)

	add_title_slide(presentation)
	add_exec_summary_slide(presentation)
	add_current_architecture_slide(presentation)
	add_mapping_slide(presentation)
	add_target_architecture_slide(presentation)
	add_benefits_slide(presentation)
	add_migration_slide(presentation)
	add_risks_slide(presentation)
	add_pilot_slide(presentation)
	add_py3270_appendix_slide(presentation)
	add_java_leanft_appendix_slide(presentation)

	output_path = Path(__file__).resolve().parent / OUTPUT_NAME
	try:
		presentation.save(output_path)
		return output_path
	except PermissionError:
		fallback_path = output_path.with_stem(f"{output_path.stem}_updated")
		presentation.save(fallback_path)
		return fallback_path


def blank_slide(prs):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	background = slide.background.fill
	background.solid()
	background.fore_color.rgb = LIGHT
	return slide


def add_footer(slide, slide_number):
	line = slide.shapes.add_shape(
		MSO_AUTO_SHAPE_TYPE.RECTANGLE,
		Inches(0),
		Inches(7.12),
		Inches(13.333),
		Inches(0.06),
	)
	line.fill.solid()
	line.fill.fore_color.rgb = BRAND
	line.line.fill.background()

	footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.17), Inches(10.8), Inches(0.22))
	frame = footer.text_frame
	frame.clear()
	paragraph = frame.paragraphs[0]
	paragraph.text = "Ascendion | BOA Jarvis to LeanFT JavaScript Migration Plan"
	paragraph.font.size = Pt(9)
	paragraph.font.color.rgb = SLATE

	number = slide.shapes.add_textbox(Inches(12.5), Inches(7.15), Inches(0.4), Inches(0.22))
	number_frame = number.text_frame
	number_frame.clear()
	paragraph = number_frame.paragraphs[0]
	paragraph.text = str(slide_number)
	paragraph.alignment = PP_ALIGN.RIGHT
	paragraph.font.size = Pt(9)
	paragraph.font.color.rgb = SLATE


def add_title(slide, title, subtitle=None):
	title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(9.6), Inches(0.7))
	title_frame = title_box.text_frame
	title_frame.clear()
	paragraph = title_frame.paragraphs[0]
	paragraph.text = title
	paragraph.font.size = Pt(28)
	paragraph.font.bold = True
	paragraph.font.color.rgb = INK

	accent = slide.shapes.add_shape(
		MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
		Inches(0.62),
		Inches(1.2),
		Inches(1.15),
		Inches(0.12),
	)
	accent.fill.solid()
	accent.fill.fore_color.rgb = BRAND
	accent.line.fill.background()

	if subtitle:
		subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(11.8), Inches(0.42))
		subtitle_frame = subtitle_box.text_frame
		subtitle_frame.clear()
		paragraph = subtitle_frame.paragraphs[0]
		paragraph.text = subtitle
		paragraph.font.size = Pt(13)
		paragraph.font.color.rgb = SLATE


def add_textbox(slide, left, top, width, height, text, font_size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
	box = slide.shapes.add_textbox(left, top, width, height)
	frame = box.text_frame
	frame.clear()
	frame.word_wrap = True
	frame.vertical_anchor = MSO_ANCHOR.TOP
	paragraph = frame.paragraphs[0]
	paragraph.text = text
	paragraph.alignment = align
	paragraph.font.size = Pt(font_size)
	paragraph.font.bold = bold
	paragraph.font.color.rgb = color
	return box


def add_bullets(slide, left, top, width, height, items, font_size=16, color=INK):
	box = slide.shapes.add_textbox(left, top, width, height)
	frame = box.text_frame
	frame.clear()
	frame.word_wrap = True
	for index, item in enumerate(items):
		paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
		paragraph.text = item
		paragraph.level = 0
		paragraph.bullet = True
		paragraph.space_after = Pt(8)
		paragraph.font.size = Pt(font_size)
		paragraph.font.color.rgb = color
	return box


def add_card(slide, left, top, width, height, heading, lines, fill_color, title_color=INK):
	card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
	card.fill.solid()
	card.fill.fore_color.rgb = fill_color
	card.line.color.rgb = WHITE

	add_textbox(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.28), heading, 14, title_color, True)

	body = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.38), width - Inches(0.3), height - Inches(0.48))
	frame = body.text_frame
	frame.clear()
	frame.word_wrap = True
	for index, line in enumerate(lines):
		paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
		paragraph.text = line
		paragraph.font.size = Pt(11)
		paragraph.font.color.rgb = INK
		paragraph.space_after = Pt(3)
	return card


def connect_vertical(slide, x_center, top, bottom):
	connector = slide.shapes.add_connector(
		MSO_CONNECTOR.STRAIGHT,
		x_center,
		top,
		x_center,
		bottom,
	)
	connector.line.color.rgb = SLATE
	connector.line.width = Pt(1.5)


def add_title_slide(prs):
	slide = blank_slide(prs)

	banner = slide.shapes.add_shape(
		MSO_AUTO_SHAPE_TYPE.RECTANGLE,
		Inches(0),
		Inches(0),
		Inches(13.333),
		Inches(0.42),
	)
	banner.fill.solid()
	banner.fill.fore_color.rgb = BRAND
	banner.line.fill.background()

	add_textbox(slide, Inches(0.7), Inches(1.0), Inches(8.6), Inches(0.9), "BOA Jarvis to LeanFT JavaScript Migration Plan", 28, INK, True)
	add_textbox(slide, Inches(0.72), Inches(1.88), Inches(8.3), Inches(0.7), "Modernizing UFT VBScript mainframe automation to a VS Code, Jasmine, and LeanFT engineering model", 16, SLATE)

	highlight = slide.shapes.add_shape(
		MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
		Inches(8.9),
		Inches(1.0),
		Inches(3.7),
		Inches(2.1),
	)
	highlight.fill.solid()
	highlight.fill.fore_color.rgb = PALE_ORANGE
	highlight.line.color.rgb = BRAND

	add_textbox(slide, Inches(9.15), Inches(1.25), Inches(3.2), Inches(0.35), "Presentation objective", 16, BRAND_DARK, True)
	add_bullets(
		slide,
		Inches(9.15),
		Inches(1.7),
		Inches(3.1),
		Inches(1.1),
		[
			"Preserve the layered Jarvis design",
			"Modernize language, tooling, and CI/CD",
			"Show a practical migration factory approach",
		],
		12,
		INK,
	)

	add_card(
		slide,
		Inches(0.75),
		Inches(3.2),
		Inches(3.65),
		Inches(2.0),
		"What this POC proves",
		[
			"LeanFT JavaScript can drive TN3270 through an existing HLLAPI-capable emulator.",
			"Jarvis concepts map cleanly into driver, keyword, object repository, data, config, and reporting layers.",
		],
		PALE_BLUE,
		BLUE,
	)
	add_card(
		slide,
		Inches(4.55),
		Inches(3.2),
		Inches(3.65),
		Inches(2.0),
		"POC flows in scope",
		[
			"Login",
			"Account Inquiry",
			"Funds Transfer",
			"Jenkins execution across SIT and UAT",
		],
		PALE_GREEN,
		GREEN,
	)
	add_card(
		slide,
		Inches(8.35),
		Inches(3.2),
		Inches(3.9),
		Inches(2.0),
		"Evidence used for this deck",
		[
			"Architecture docs in docs/01 and docs/02",
			"Framework context files",
			"Actual runner.js, terminalHelper.js, screens.js, and Jasmine spec implementation",
		],
		PALE_GOLD,
		AMBER,
	)

	add_textbox(slide, Inches(0.75), Inches(5.7), Inches(4.0), Inches(0.3), f"Prepared on {date.today().strftime('%d %b %Y')}", 11, SLATE)
	add_footer(slide, 1)


def add_exec_summary_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Executive Summary", "The migration does not require rethinking the testing model. It modernizes the implementation stack and delivery flow.")

	add_card(
		slide,
		Inches(0.7),
		Inches(1.9),
		Inches(4.0),
		Inches(1.7),
		"Current-state constraint",
		[
			"Jarvis depends on VBScript, UFT artifacts, XML flows, and binary object repository assets that are harder to diff, review, and scale.",
		],
		PALE_RED,
		RED,
	)
	add_card(
		slide,
		Inches(4.95),
		Inches(1.9),
		Inches(4.0),
		Inches(1.7),
		"Target-state shift",
		[
			"The new framework keeps the same layered design but uses LeanFT JavaScript, Jasmine, plain-text screen definitions, JSON or Excel-backed test data, and Jenkins-ready execution.",
		],
		PALE_BLUE,
		BLUE,
	)
	add_card(
		slide,
		Inches(9.2),
		Inches(1.9),
		Inches(3.45),
		Inches(1.7),
		"Business outcome",
		[
			"Faster onboarding, cleaner Git visibility, better IDE support, and an AI-assisted migration pattern.",
		],
		PALE_GREEN,
		GREEN,
	)

	add_textbox(slide, Inches(0.8), Inches(4.0), Inches(3.5), Inches(0.28), "What remains familiar", 16, INK, True)
	add_bullets(
		slide,
		Inches(0.8),
		Inches(4.35),
		Inches(5.4),
		Inches(1.8),
		[
			"Layered automation design",
			"Keyword-driven mainframe interactions",
			"Central object repository and external configuration",
			"Data-driven execution and Jenkins orchestration",
		],
		15,
	)

	add_textbox(slide, Inches(6.8), Inches(4.0), Inches(3.5), Inches(0.28), "What improves materially", 16, INK, True)
	add_bullets(
		slide,
		Inches(6.8),
		Inches(4.35),
		Inches(5.5),
		Inches(2.0),
		[
			"VBScript to modern JavaScript",
			"Binary assets to plain-text source-controlled files",
			"UFT-centric authoring to VS Code + GitHub Copilot",
			"Manual migration effort reduced by a deterministic BOA migration agent",
		],
		15,
	)

	add_footer(slide, 2)


def add_current_architecture_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Current State Architecture: Jarvis / UFT / VBScript", "The current model is modular, but several layers are implemented with older and more opaque technologies.")

	cards = [
		(Inches(0.8), Inches(2.0), "Test Suite", ["Spreadsheet-driven suite definitions", "Execution entry point for business scenarios"], PALE_GOLD),
		(Inches(4.55), Inches(2.0), "Test Flows", ["XML flow orchestration", "Externalized step sequencing"], PALE_GOLD),
		(Inches(8.3), Inches(2.0), "Test Data", ["Spreadsheet-based data sets", "Separate from scripts"], PALE_GOLD),
		(Inches(0.8), Inches(3.7), "Execution Engine", ["Jarvis / UFT driver-runner model", "Scenario coordination and playback"], PALE_RED),
		(Inches(4.55), Inches(3.7), "Libraries", [".qfl VBScript keyword libraries", "Reusable screen actions"], PALE_RED),
		(Inches(8.3), Inches(3.7), "Object Repository", ["UFT .tsr or .mtr assets", "Binary-style repository management"], PALE_RED),
		(Inches(0.8), Inches(5.4), "Configuration", ["Config templates and runtime settings"], PALE_GRAY),
		(Inches(4.55), Inches(5.4), "Reporting", ["UFT HTML reports and zipped results"], PALE_GRAY),
		(Inches(8.3), Inches(5.4), "CI/CD", ["Jenkins integration present, but framework tooling remains UFT-centric"], PALE_GRAY),
	]

	for left, top, heading, lines, fill in cards:
		add_card(slide, left, top, Inches(3.0), Inches(1.15), heading, lines, fill)

	note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.75), Inches(1.1), Inches(1.85), Inches(0.55))
	note.fill.solid()
	note.fill.fore_color.rgb = BRAND
	note.line.fill.background()
	add_textbox(slide, Inches(10.92), Inches(1.23), Inches(1.5), Inches(0.2), "Pain points", 14, WHITE, True, PP_ALIGN.CENTER)
	add_bullets(
		slide,
		Inches(10.35),
		Inches(1.75),
		Inches(2.0),
		Inches(3.0),
		[
			"VBScript deprecation risk",
			"Harder Git diff and code review experience",
			"Less developer-friendly tooling",
			"Manual migration of legacy assets is slower without structured accelerators",
		],
		12,
	)

	add_footer(slide, 3)


def add_mapping_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Migration Mapping: Jarvis Components to the New Framework", "The POC preserves the business testing intent while swapping each legacy layer for a modern equivalent.")

	data = [
		["Jarvis component", "Current technology", "Target component", "Target technology"],
		["Test Suite", "Excel spreadsheet", "spec/ftdtest_jasmine_spec.js", "Jasmine JavaScript spec"],
		["Test Flow", "XML", "describe / it blocks", "Native JavaScript structure"],
		["Test Data", "Excel spreadsheet", "testdata JSON + xlsx support", "JSON or Excel via Node.js"],
		["Execution Engine", "UFT proprietary runner", "driver/runner.js", "LeanFT SDK lifecycle"],
		["Libraries", ".qfl VBScript", "libraries/terminalHelper.js", "Reusable JS keywords"],
		["Object Repository", "UFT .tsr / .mtr", "objectrepository/screens.js", "Plain-text screen metadata"],
		["Configuration", "Config templates", "config/settings.js", "Node module + env vars"],
		["Reporting", "UFT report", "results/ LeanFT report", "HTML report with snapshots"],
		["CI/CD", "Jenkinsfile", "Jenkinsfile", "Parameterized SIT / UAT execution"],
	]

	table = slide.shapes.add_table(len(data), len(data[0]), Inches(0.55), Inches(1.8), Inches(12.2), Inches(4.7)).table
	widths = [2.2, 2.2, 3.0, 4.8]
	for index, width in enumerate(widths):
		table.columns[index].width = Inches(width)

	for row_index, row in enumerate(data):
		for col_index, value in enumerate(row):
			cell = table.cell(row_index, col_index)
			cell.text = value
			paragraph = cell.text_frame.paragraphs[0]
			paragraph.font.size = Pt(11 if row_index else 12)
			paragraph.font.bold = row_index == 0
			paragraph.font.color.rgb = INK
			cell.fill.solid()
			if row_index == 0:
				cell.fill.fore_color.rgb = PALE_ORANGE
			elif row_index % 2:
				cell.fill.fore_color.rgb = WHITE
			else:
				cell.fill.fore_color.rgb = PALE_GRAY

	add_textbox(slide, Inches(0.7), Inches(6.72), Inches(12.0), Inches(0.28), "Key client message: this is a migration with continuity, not a rewrite with loss of control.", 13, BRAND_DARK, True)
	add_footer(slide, 4)


def add_target_architecture_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Target Architecture: LeanFT JavaScript Framework", "The implementation already present in the POC follows a modern, reviewable, and modular structure.")

	x_positions = [Inches(0.8), Inches(3.95), Inches(7.1), Inches(10.25)]
	top_positions = [Inches(2.0), Inches(4.05)]
	boxes = [
		(0, 0, "Test Suite", ["Jasmine spec", "Business scenarios and assertions"], PALE_BLUE),
		(1, 0, "Driver / Runner", ["LeanFT init / cleanup", "Shared TE session lifecycle"], PALE_ORANGE),
		(2, 0, "Action Layer", ["terminalHelper.js keywords", "Mainframe business actions"], PALE_GREEN),
		(3, 0, "Object Repository", ["screens.js", "Identifiers, fields, and key maps"], PALE_GOLD),
		(0, 1, "Test Data", ["JSON default", "Excel support via xlsx"], PALE_GRAY),
		(1, 1, "Configuration", ["settings.js", "TE short name, env, timeouts"], PALE_GRAY),
		(2, 1, "Reporting", ["LeanFT HTML report", "Snapshots and pass/fail logging"], PALE_GRAY),
		(3, 1, "CI/CD", ["Jenkins SIT / UAT", "Filterable suite execution"], PALE_GRAY),
	]

	for col, row, heading, lines, fill in boxes:
		add_card(slide, x_positions[col], top_positions[row], Inches(2.65), Inches(1.25), heading, lines, fill)

	app = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(5.8), Inches(4.9), Inches(0.78))
	app.fill.solid()
	app.fill.fore_color.rgb = RGBColor(226, 232, 240)
	app.line.color.rgb = SLATE
	add_textbox(slide, Inches(4.45), Inches(5.98), Inches(4.4), Inches(0.22), "IBM Mainframe through an existing HLLAPI-capable terminal emulator", 14, INK, True, PP_ALIGN.CENTER)

	connect_vertical(slide, Inches(5.28), Inches(3.26), Inches(3.95))
	connect_vertical(slide, Inches(8.43), Inches(3.26), Inches(3.95))
	connect_vertical(slide, Inches(8.43), Inches(5.3), Inches(5.8))

	add_textbox(slide, Inches(0.82), Inches(6.7), Inches(12.0), Inches(0.28), "Important architecture note: LeanFT does not open TN3270 directly; it automates an already-running emulator session via HLLAPI.", 12, SLATE)
	add_footer(slide, 5)


def add_benefits_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Why This Framework Is Better for Delivery", "The technical modernization is paired with a better authoring and migration experience for the QE team.")

	add_textbox(slide, Inches(0.8), Inches(1.95), Inches(4.7), Inches(0.3), "Engineering benefits", 17, INK, True)
	add_bullets(
		slide,
		Inches(0.8),
		Inches(2.3),
		Inches(5.3),
		Inches(3.4),
		[
			"JavaScript is closer to VBScript than Java or C#, which reduces training friction.",
			"VS Code provides a lightweight and familiar IDE with Git-native workflows.",
			"Plain-text repository assets are easier to review, diff, branch, and promote.",
			"Jasmine plus npm gives a simple, CI-friendly execution model.",
			"JSON-first test data can still coexist with Excel input when needed.",
		],
		15,
	)

	add_textbox(slide, Inches(6.8), Inches(1.95), Inches(4.9), Inches(0.3), "Acceleration benefits", 17, INK, True)
	add_bullets(
		slide,
		Inches(6.8),
		Inches(2.3),
		Inches(5.3),
		Inches(3.4),
		[
			"GitHub Copilot can assist with routine coding, refactoring, and test authoring inside VS Code.",
			"The repo includes a BOA migration agent that reads project context, maps VBScript constructs, generates Jasmine skeletons, and flags only the steps that need manual intervention.",
			"The agent provides a migration delta report, giving traceability from automated translation to remaining manual work.",
			"The framework is positioned for stronger engineering collaboration between QE and development teams.",
		],
		15,
	)

	benefit_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.2), Inches(11.6), Inches(0.72))
	benefit_bar.fill.solid()
	benefit_bar.fill.fore_color.rgb = PALE_ORANGE
	benefit_bar.line.color.rgb = BRAND
	add_textbox(slide, Inches(1.15), Inches(6.4), Inches(11.0), Inches(0.2), "Client takeaway: the target stack improves both automation maintainability and migration throughput.", 15, BRAND_DARK, True, PP_ALIGN.CENTER)
	add_footer(slide, 6)


def add_migration_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Suggested Migration Approach", "A factory model is safer than a big-bang conversion. The POC and BOA migration agent support this sequence.")

	phases = [
		("1. Assess", "Inventory current suites, libraries, flows, data, and object repository dependencies."),
		("2. Map", "Map Jarvis constructs to the LeanFT JavaScript framework and identify manual-only gaps."),
		("3. Baseline", "Stabilize target framework layers, CI settings, reporting, and environment conventions."),
		("4. Pilot", "Migrate representative high-value flows such as login, account inquiry, and funds transfer."),
		("5. Industrialize", "Use Copilot and the migration agent to accelerate repetitive conversions with guardrails."),
		("6. Scale", "Expand by application area, validate SIT and UAT behavior, and retire legacy suites in waves."),
	]

	left = Inches(0.75)
	top = Inches(2.0)
	width = Inches(12.0)
	gap = Inches(0.12)
	box_width = (width - gap * 5) / 6
	for index, (heading, detail) in enumerate(phases):
		x_pos = left + index * (box_width + gap)
		fill = PALE_ORANGE if index in (1, 4) else PALE_BLUE if index in (0, 3) else PALE_GREEN
		add_card(slide, x_pos, top, box_width, Inches(2.2), heading, [detail], fill)

	add_textbox(slide, Inches(0.82), Inches(4.65), Inches(3.0), Inches(0.25), "How the accelerator helps", 16, INK, True)
	add_bullets(
		slide,
		Inches(0.82),
		Inches(4.95),
		Inches(6.2),
		Inches(1.4),
		[
			"Loads framework and migration context before translation",
			"Maps known VBScript and UFT constructs deterministically",
			"Creates skeletons, keyword translations, and a migration delta report",
		],
		14,
	)

	add_textbox(slide, Inches(7.0), Inches(4.65), Inches(3.3), Inches(0.25), "Operating principle", 16, INK, True)
	add_bullets(
		slide,
		Inches(7.0),
		Inches(4.95),
		Inches(5.3),
		Inches(1.5),
		[
			"Automate the known mappings",
			"Flag anything that requires live environment confirmation",
			"Do not guess locators, credentials, or environment-specific identifiers",
		],
		14,
	)

	add_footer(slide, 7)


def add_risks_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Migration Risks and How to Control Them", "The main risks are known upfront and can be managed with the right pilot discipline.")

	headers = ["Risk or dependency", "Mitigation in the proposed plan"]
	rows = [
		["Screen locators must match live emulator labels", "Capture and validate attachedText values using the LeanFT Object Identification Center; do not guess them in code."],
		["Screen identifiers can vary by environment", "Validate identifiers in SIT and UAT during pilot onboarding and keep them centralized in screens.js."],
		["Credentials and sensitive data must not be embedded", "Keep credentials out of specs and use managed external test data and environment configuration."],
		["Some custom Jarvis keywords may not have a direct map", "Use the migration agent to flag keyword gaps and implement them once in terminalHelper.js for reuse."],
		["Citrix or non-3270 surfaces may require separate treatment", "Keep the first migration wave focused on standard TN3270 automation; assess exceptional channels separately."],
		["Execution still depends on a running emulator session", "Document emulator setup, session short names, and Jenkins runtime configuration as part of environment readiness."],
	]

	table = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.7), Inches(1.95), Inches(12.0), Inches(4.9)).table
	table.columns[0].width = Inches(3.5)
	table.columns[1].width = Inches(8.5)

	for col, header in enumerate(headers):
		cell = table.cell(0, col)
		cell.text = header
		cell.fill.solid()
		cell.fill.fore_color.rgb = PALE_ORANGE
		paragraph = cell.text_frame.paragraphs[0]
		paragraph.font.size = Pt(12)
		paragraph.font.bold = True
		paragraph.font.color.rgb = INK

	for row_index, row in enumerate(rows, start=1):
		for col_index, value in enumerate(row):
			cell = table.cell(row_index, col_index)
			cell.text = value
			cell.fill.solid()
			cell.fill.fore_color.rgb = WHITE if row_index % 2 else PALE_GRAY
			paragraph = cell.text_frame.paragraphs[0]
			paragraph.font.size = Pt(11)
			paragraph.font.color.rgb = INK

	add_footer(slide, 8)


def add_pilot_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Recommended Pilot and Next Steps", "Use the current POC as the demonstration baseline, then move into a controlled pilot migration wave.")

	add_card(
		slide,
		Inches(0.8),
		Inches(2.0),
		Inches(3.8),
		Inches(2.15),
		"Pilot scope",
		[
			"3 to 5 representative business flows",
			"At least one login or navigation pattern",
			"At least one inquiry flow",
			"At least one transaction flow",
			"SIT and UAT execution through Jenkins",
		],
		PALE_BLUE,
		BLUE,
	)
	add_card(
		slide,
		Inches(4.78),
		Inches(2.0),
		Inches(3.8),
		Inches(2.15),
		"Success criteria",
		[
			"Functional parity with selected Jarvis flows",
			"Stable object repository and environment settings",
			"Repeatable CI execution and report publication",
			"Measured reduction in migration effort for each additional script",
		],
		PALE_GREEN,
		GREEN,
	)
	add_card(
		slide,
		Inches(8.76),
		Inches(2.0),
		Inches(3.8),
		Inches(2.15),
		"Immediate next steps",
		[
			"Confirm pilot scope and environments",
			"Capture live locators for pilot screens",
			"Prioritize Jarvis suites by reuse and business value",
			"Define governance for AI-assisted migration review",
		],
		PALE_ORANGE,
		BRAND_DARK,
	)

	add_card(
		slide,
		Inches(0.85),
		Inches(4.65),
		Inches(5.85),
		Inches(1.75),
		"What we need from Bank of America to continue",
		[
			"Access to a BOA test environment with the actual mainframe application and representative user journeys.",
			"One BOA-side SME, tester, or engineer working closely with us to validate screens, navigation, and expected results.",
			"Support to configure the terminal emulator, session short names, host connectivity, and any environment-specific runtime settings.",
		],
		PALE_RED,
		RED,
	)

	add_card(
		slide,
		Inches(6.95),
		Inches(4.65),
		Inches(5.5),
		Inches(1.75),
		"Dependencies and open questions for next steps",
		[
			"Which BOA environments are available for SIT, UAT, or pilot execution?",
			"What emulator, access controls, credentials handling, and network prerequisites are required?",
			"Which flows, test data, and application variants should be the first migration wave?",
			"Are there any BOA-specific dependencies such as Citrix, middleware hops, or security approvals we should plan for?",
		],
		PALE_GOLD,
		AMBER,
	)

	quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.48), Inches(11.2), Inches(0.42))
	quote.fill.solid()
	quote.fill.fore_color.rgb = PALE_ORANGE
	quote.line.color.rgb = BRAND
	add_textbox(slide, Inches(1.2), Inches(6.58), Inches(10.8), Inches(0.18), "Client ask: to move beyond the current POC, one of the team members should work with BOA in the real environment to complete setup, validate locators, and industrialize the migration path.", 12, BRAND_DARK, True, PP_ALIGN.CENTER)

	add_footer(slide, 9)


def add_py3270_appendix_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Appendix: Python / py3270 Open-Source POC", "A second branch demonstrates a lower-cost modernization path using Python, pytest, and direct TN3270 connectivity.")

	add_card(
		slide,
		Inches(0.8),
		Inches(2.0),
		Inches(3.8),
		Inches(2.2),
		"What this POC is",
		[
			"Branch: python_based_open_source_poc",
			"Stack: Python 3.10+, py3270, pytest, pytest-html",
			"Direct TN3270 over s3270 subprocess, no LeanFT, no HLLAPI, no emulator license",
		],
		PALE_BLUE,
		BLUE,
	)

	add_card(
		slide,
		Inches(4.78),
		Inches(2.0),
		Inches(3.8),
		Inches(2.2),
		"Architecture in one line",
		[
			"pytest suite -> runner.py -> terminal_helper.py -> py3270 -> s3270 -> TN3270 socket -> mainframe",
			"The same layered model is retained: config, driver, action layer, object repository, test data, reporting, and Jenkins CI",
		],
		PALE_GREEN,
		GREEN,
	)

	add_card(
		slide,
		Inches(8.76),
		Inches(2.0),
		Inches(3.8),
		Inches(2.2),
		"Why mention it to the client",
		[
			"Zero licensing cost",
			"Headless execution model suited for CI",
			"Useful as an alternate modernization option if BOA prefers open-source tooling",
		],
		PALE_ORANGE,
		BRAND_DARK,
	)

	add_textbox(slide, Inches(0.82), Inches(4.7), Inches(3.8), Inches(0.25), "Key differences from the LeanFT POC", 16, INK, True)
	add_bullets(
		slide,
		Inches(0.82),
		Inches(5.0),
		Inches(5.7),
		Inches(1.35),
		[
			"Python instead of JavaScript",
			"pytest instead of Jasmine",
			"Direct socket connectivity instead of emulator-driven HLLAPI automation",
		],
		14,
	)

	callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.95), Inches(5.45), Inches(1.1))
	callout.fill.solid()
	callout.fill.fore_color.rgb = PALE_GOLD
	callout.line.color.rgb = AMBER
	add_textbox(slide, Inches(7.05), Inches(5.22), Inches(4.95), Inches(0.45), "Suggested positioning: present this as an optional open-source alternate path, not a replacement for the primary LeanFT migration story unless BOA wants a no-license model.", 13, BRAND_DARK, True, PP_ALIGN.CENTER)

	add_footer(slide, 10)


def add_java_leanft_appendix_slide(prs):
	slide = blank_slide(prs)
	add_title(slide, "Appendix: Java / LeanFT POC", "This branch also shows a Java-based LeanFT implementation for teams that prefer a Java and Maven stack.")

	add_card(
		slide,
		Inches(0.8),
		Inches(2.0),
		Inches(3.8),
		Inches(2.2),
		"What this POC is",
		[
			"Stack: Java 11+, LeanFT runtime engine, TestNG, Maven",
			"Main files: BaseTest.java, TerminalHelper.java, Screens.java, FtdMainframeTest.java",
			"LeanFT HTML reporting with screenshots and step logs",
		],
		PALE_BLUE,
		BLUE,
	)

	add_card(
		slide,
		Inches(4.78),
		Inches(2.0),
		Inches(3.8),
		Inches(2.2),
		"Architecture in one line",
		[
			"TestNG suite -> BaseTest -> TerminalHelper -> LeanFT TE SDK -> terminal emulator session -> mainframe",
			"The same layered model is preserved with config, data provider, object repository, action layer, reporting, and Jenkins CI",
		],
		PALE_GREEN,
		GREEN,
	)

	add_card(
		slide,
		Inches(8.76),
		Inches(2.0),
		Inches(3.8),
		Inches(2.2),
		"Why mention it to the client",
		[
			"Suitable if BOA has stronger Java engineering alignment",
			"Fits enterprise Maven and TestNG delivery patterns",
			"Still modernizes Jarvis while staying within the LeanFT ecosystem",
		],
		PALE_ORANGE,
		BRAND_DARK,
	)

	add_textbox(slide, Inches(0.82), Inches(4.7), Inches(4.1), Inches(0.25), "Key characteristics", 16, INK, True)
	add_bullets(
		slide,
		Inches(0.82),
		Inches(5.0),
		Inches(5.8),
		Inches(1.35),
		[
			"Java instead of JavaScript or Python",
			"TestNG and Maven instead of Jasmine or pytest",
			"Requires LeanFT runtime and terminal emulator-based execution",
		],
		14,
	)

	callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.95), Inches(5.45), Inches(1.1))
	callout.fill.solid()
	callout.fill.fore_color.rgb = PALE_GOLD
	callout.line.color.rgb = AMBER
	add_textbox(slide, Inches(7.05), Inches(5.18), Inches(4.95), Inches(0.5), "Suggested positioning: present this as a second LeanFT-based option for teams that want a typed Java stack and Maven-style enterprise delivery rather than the JavaScript implementation.", 13, BRAND_DARK, True, PP_ALIGN.CENTER)

	add_footer(slide, 11)


if __name__ == "__main__":
	path = build_deck()
	print(f"Created {path}")
